"""
File Creator Service - Extended Edition
Erstellt PPTX, DOCX, XLSX und PDF Dateien mit professionellem Design
"""

import os
import uuid
import json
import shutil
from datetime import datetime, timedelta
from pathlib import Path
from typing import Optional, List, Dict, Any
from contextlib import asynccontextmanager
import asyncio

from fastapi import FastAPI, HTTPException, UploadFile, File, Form
from fastapi.responses import FileResponse, JSONResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field

# Document Libraries
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE

from openpyxl import Workbook, load_workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.chart import BarChart, LineChart, PieChart, Reference
from openpyxl.utils import get_column_letter
from openpyxl.drawing.image import Image as XLImage

from reportlab.lib.pagesizes import letter, A4
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch, cm
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image as RLImage, PageBreak
from reportlab.pdfgen import canvas

# ============================================================================
# CONFIGURATION
# ============================================================================

PORT = int(os.getenv("PORT", 8002))
BASE_URL = os.getenv("BASE_URL", f"http://localhost:{PORT}")
FILES_DIR = Path(os.getenv("FILES_DIR", "/app/files"))
ASSETS_DIR = Path(os.getenv("ASSETS_DIR", "/app/assets"))
TEMPLATES_DIR = Path(os.getenv("TEMPLATES_DIR", "/app/templates"))
FILE_RETENTION_HOURS = int(os.getenv("FILE_RETENTION_HOURS", 24))

# Ensure directories exist
FILES_DIR.mkdir(parents=True, exist_ok=True)
ASSETS_DIR.mkdir(parents=True, exist_ok=True)
TEMPLATES_DIR.mkdir(parents=True, exist_ok=True)

# Subdirectories for assets
(ASSETS_DIR / "logos").mkdir(exist_ok=True)
(ASSETS_DIR / "images").mkdir(exist_ok=True)

# ============================================================================
# COLOR PALETTES FOR PROFESSIONAL DESIGN
# ============================================================================

COLOR_PALETTES = {
    "midnight_executive": {
        "primary": "1E2761",
        "secondary": "CADCFC", 
        "accent": "FFFFFF",
        "text_dark": "1E2761",
        "text_light": "FFFFFF"
    },
    "forest_moss": {
        "primary": "2C5F2D",
        "secondary": "97BC62",
        "accent": "F5F5F5",
        "text_dark": "2C5F2D",
        "text_light": "FFFFFF"
    },
    "coral_energy": {
        "primary": "F96167",
        "secondary": "F9E795",
        "accent": "2F3C7E",
        "text_dark": "2F3C7E",
        "text_light": "FFFFFF"
    },
    "warm_terracotta": {
        "primary": "B85042",
        "secondary": "E7E8D1",
        "accent": "A7BEAE",
        "text_dark": "B85042",
        "text_light": "FFFFFF"
    },
    "ocean_gradient": {
        "primary": "065A82",
        "secondary": "1C7293",
        "accent": "21295C",
        "text_dark": "21295C",
        "text_light": "FFFFFF"
    },
    "charcoal_minimal": {
        "primary": "36454F",
        "secondary": "F2F2F2",
        "accent": "212121",
        "text_dark": "212121",
        "text_light": "FFFFFF"
    },
    "teal_trust": {
        "primary": "028090",
        "secondary": "00A896",
        "accent": "02C39A",
        "text_dark": "028090",
        "text_light": "FFFFFF"
    },
    "berry_cream": {
        "primary": "6D2E46",
        "secondary": "A26769",
        "accent": "ECE2D0",
        "text_dark": "6D2E46",
        "text_light": "FFFFFF"
    },
    "corporate_blue": {
        "primary": "003366",
        "secondary": "0066CC",
        "accent": "FF9900",
        "text_dark": "003366",
        "text_light": "FFFFFF"
    }
}

# ============================================================================
# PYDANTIC MODELS
# ============================================================================

# --- PowerPoint Models ---
class SlideContent(BaseModel):
    title: str
    content: Optional[str] = None
    bullet_points: Optional[List[str]] = None
    layout: str = "title_content"  # title_only, title_content, two_column, image_left, image_right, stats, comparison
    image_path: Optional[str] = None
    stats: Optional[List[Dict[str, str]]] = None  # [{"value": "95%", "label": "Kundenzufriedenheit"}]
    columns: Optional[List[Dict[str, Any]]] = None  # Für two_column layout

class PresentationRequest(BaseModel):
    title: str
    subtitle: Optional[str] = None
    author: Optional[str] = None
    slides: List[SlideContent]
    color_palette: str = "corporate_blue"
    logo: Optional[str] = None  # Logo-Name aus assets/logos
    include_toc: bool = False
    include_closing: bool = True

# --- Word Document Models ---
class DocumentSection(BaseModel):
    heading: Optional[str] = None
    heading_level: int = 1
    content: Optional[str] = None
    bullet_points: Optional[List[str]] = None
    table: Optional[List[List[str]]] = None
    image_path: Optional[str] = None

class WordDocumentRequest(BaseModel):
    title: str
    subtitle: Optional[str] = None
    author: Optional[str] = None
    sections: List[DocumentSection]
    template: Optional[str] = None  # Template-Name aus templates/
    logo: Optional[str] = None
    header_text: Optional[str] = None
    footer_text: Optional[str] = None
    include_toc: bool = False

# --- Excel Models ---
class ChartConfig(BaseModel):
    type: str = "bar"  # bar, line, pie
    title: str
    data_range: str  # z.B. "A1:D5"
    categories_range: Optional[str] = None
    position: str = "E2"

class ExcelSheet(BaseModel):
    name: str
    headers: Optional[List[str]] = None
    data: Optional[List[List[Any]]] = None
    column_widths: Optional[Dict[str, int]] = None
    charts: Optional[List[ChartConfig]] = None
    formulas: Optional[Dict[str, str]] = None  # {"C10": "=SUM(C2:C9)"}
    conditional_formatting: Optional[Dict[str, Any]] = None

class ExcelRequest(BaseModel):
    filename: str
    sheets: List[ExcelSheet]
    template: Optional[str] = None
    logo: Optional[str] = None
    style: str = "professional"  # professional, financial, minimal

# --- PDF Models ---
class PDFSection(BaseModel):
    type: str = "paragraph"  # paragraph, heading, table, image, spacer, page_break
    content: Optional[str] = None
    level: int = 1  # Für headings
    data: Optional[List[List[str]]] = None  # Für tables
    image_path: Optional[str] = None
    height: Optional[float] = None  # Spacer height in cm

class PDFRequest(BaseModel):
    title: str
    author: Optional[str] = None
    sections: List[PDFSection]
    page_size: str = "A4"  # A4, letter
    logo: Optional[str] = None
    header_text: Optional[str] = None
    footer_text: Optional[str] = None
    color_scheme: str = "corporate_blue"

# --- Asset/Logo Models ---
class LogoUploadResponse(BaseModel):
    success: bool
    logo_name: str
    message: str

class TemplateUploadResponse(BaseModel):
    success: bool
    template_name: str
    template_type: str
    message: str

# ============================================================================
# HELPER FUNCTIONS
# ============================================================================

def generate_filename(prefix: str, extension: str) -> str:
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    unique_id = uuid.uuid4().hex[:8]
    return f"{prefix}_{timestamp}_{unique_id}.{extension}"

def hex_to_rgb(hex_color: str) -> tuple:
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def get_logo_path(logo_name: Optional[str]) -> Optional[Path]:
    if not logo_name:
        return None
    logo_path = ASSETS_DIR / "logos" / logo_name
    if logo_path.exists():
        return logo_path
    # Versuche mit verschiedenen Erweiterungen
    for ext in ['.png', '.jpg', '.jpeg', '.svg']:
        test_path = ASSETS_DIR / "logos" / f"{logo_name}{ext}"
        if test_path.exists():
            return test_path
    return None

def get_template_path(template_name: Optional[str], template_type: str) -> Optional[Path]:
    if not template_name:
        return None
    template_path = TEMPLATES_DIR / template_type / template_name
    if template_path.exists():
        return template_path
    return None

async def cleanup_old_files():
    """Löscht Dateien älter als FILE_RETENTION_HOURS"""
    while True:
        cutoff = datetime.now() - timedelta(hours=FILE_RETENTION_HOURS)
        for file_path in FILES_DIR.iterdir():
            if file_path.is_file():
                file_time = datetime.fromtimestamp(file_path.stat().st_mtime)
                if file_time < cutoff:
                    file_path.unlink()
        await asyncio.sleep(3600)  # Check every hour

# ============================================================================
# POWERPOINT CREATION (PROFESSIONAL DESIGN)
# ============================================================================

def create_professional_pptx(request: PresentationRequest) -> Path:
    """Erstellt eine professionell gestaltete PowerPoint-Präsentation"""
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    palette = COLOR_PALETTES.get(request.color_palette, COLOR_PALETTES["corporate_blue"])
    logo_path = get_logo_path(request.logo)
    
    # --- Title Slide ---
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Dark background for title
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RgbColor(*hex_to_rgb(palette["primary"]))
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.5), Inches(11.8), Inches(1.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = request.title
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RgbColor(*hex_to_rgb(palette["text_light"]))
    title_para.alignment = PP_ALIGN.LEFT
    
    # Subtitle
    if request.subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.75), Inches(4.2), Inches(11.8), Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = request.subtitle
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = RgbColor(*hex_to_rgb(palette["secondary"]))
        subtitle_para.alignment = PP_ALIGN.LEFT
    
    # Logo on title slide
    if logo_path:
        slide.shapes.add_picture(str(logo_path), Inches(10.5), Inches(0.5), width=Inches(2))
    
    # Author and Date
    if request.author:
        author_box = slide.shapes.add_textbox(Inches(0.75), Inches(6.5), Inches(6), Inches(0.5))
        author_frame = author_box.text_frame
        author_para = author_frame.paragraphs[0]
        author_para.text = f"{request.author} | {datetime.now().strftime('%B %Y')}"
        author_para.font.size = Pt(14)
        author_para.font.color.rgb = RgbColor(*hex_to_rgb(palette["accent"]))
    
    # --- Content Slides ---
    for slide_content in request.slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Light background for content slides
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RgbColor(255, 255, 255)
        bg.line.fill.background()
        
        # Accent bar at top
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.1)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = RgbColor(*hex_to_rgb(palette["primary"]))
        accent_bar.line.fill.background()
        
        # Slide title
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.8), Inches(0.9))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = slide_content.title
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = RgbColor(*hex_to_rgb(palette["text_dark"]))
        
        # Content based on layout
        if slide_content.layout == "title_content":
            _add_standard_content(slide, slide_content, palette)
        elif slide_content.layout == "two_column":
            _add_two_column_content(slide, slide_content, palette)
        elif slide_content.layout == "stats":
            _add_stats_content(slide, slide_content, palette)
        elif slide_content.layout == "image_left":
            _add_image_left_content(slide, slide_content, palette)
        elif slide_content.layout == "image_right":
            _add_image_right_content(slide, slide_content, palette)
        else:
            _add_standard_content(slide, slide_content, palette)
        
        # Logo on each slide (small, bottom right)
        if logo_path:
            slide.shapes.add_picture(str(logo_path), Inches(11.5), Inches(6.8), width=Inches(1.2))
    
    # --- Closing Slide ---
    if request.include_closing:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RgbColor(*hex_to_rgb(palette["primary"]))
        bg.line.fill.background()
        
        thanks_box = slide.shapes.add_textbox(Inches(0.75), Inches(3), Inches(11.8), Inches(1.5))
        thanks_frame = thanks_box.text_frame
        thanks_para = thanks_frame.paragraphs[0]
        thanks_para.text = "Vielen Dank!"
        thanks_para.font.size = Pt(54)
        thanks_para.font.bold = True
        thanks_para.font.color.rgb = RgbColor(*hex_to_rgb(palette["text_light"]))
        thanks_para.alignment = PP_ALIGN.CENTER
        
        if logo_path:
            slide.shapes.add_picture(str(logo_path), Inches(5.5), Inches(5), width=Inches(2.5))
    
    # Save
    filename = generate_filename("presentation", "pptx")
    filepath = FILES_DIR / filename
    prs.save(filepath)
    return filepath

def _add_standard_content(slide, content: SlideContent, palette: dict):
    """Standard Layout: Titel + Content/Bullets"""
    y_start = Inches(1.6)
    
    if content.content:
        content_box = slide.shapes.add_textbox(Inches(0.75), y_start, Inches(11.8), Inches(4.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.text = content.content
        para.font.size = Pt(18)
        para.font.color.rgb = RgbColor(*hex_to_rgb(palette["text_dark"]))
        para.space_after = Pt(12)
    
    if content.bullet_points:
        bullet_y = y_start if not content.content else Inches(3)
        for i, point in enumerate(content.bullet_points):
            bullet_box = slide.shapes.add_textbox(Inches(0.75), bullet_y + Inches(i * 0.7), Inches(11.8), Inches(0.6))
            tf = bullet_box.text_frame
            para = tf.paragraphs[0]
            
            # Icon circle
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(0.75), bullet_y + Inches(i * 0.7) + Inches(0.1), Inches(0.25), Inches(0.25)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = RgbColor(*hex_to_rgb(palette["secondary"]))
            circle.line.fill.background()
            
            para.text = f"    {point}"
            para.font.size = Pt(16)
            para.font.color.rgb = RgbColor(*hex_to_rgb(palette["text_dark"]))

def _add_two_column_content(slide, content: SlideContent, palette: dict):
    """Zwei-Spalten Layout"""
    if content.columns and len(content.columns) >= 2:
        left_col = content.columns[0]
        right_col = content.columns[1]
        
        # Left column
        left_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(5.5), Inches(5))
        _fill_column(left_box, left_col, palette)
        
        # Right column
        right_box = slide.shapes.add_textbox(Inches(7), Inches(1.6), Inches(5.5), Inches(5))
        _fill_column(right_box, right_col, palette)

def _fill_column(textbox, col_data: dict, palette: dict):
    tf = textbox.text_frame
    tf.word_wrap = True
    if col_data.get("title"):
        para = tf.paragraphs[0]
        para.text = col_data["title"]
        para.font.size = Pt(20)
        para.font.bold = True
        para.font.color.rgb = RgbColor(*hex_to_rgb(palette["primary"]))
    if col_data.get("content"):
        para = tf.add_paragraph()
        para.text = col_data["content"]
        para.font.size = Pt(14)
        para.font.color.rgb = RgbColor(*hex_to_rgb(palette["text_dark"]))

def _add_stats_content(slide, content: SlideContent, palette: dict):
    """Stats Layout: Große Zahlen mit Labels"""
    if not content.stats:
        return
    
    num_stats = len(content.stats)
    start_x = 0.75
    stat_width = (12 - 0.75) / num_stats
    
    for i, stat in enumerate(content.stats):
        x = Inches(start_x + i * stat_width)
        
        # Value (large)
        value_box = slide.shapes.add_textbox(x, Inches(2.5), Inches(stat_width - 0.3), Inches(1.5))
        tf = value_box.text_frame
        para = tf.paragraphs[0]
        para.text = stat.get("value", "")
        para.font.size = Pt(60)
        para.font.bold = True
        para.font.color.rgb = RgbColor(*hex_to_rgb(palette["primary"]))
        para.alignment = PP_ALIGN.CENTER
        
        # Label
        label_box = slide.shapes.add_textbox(x, Inches(4.2), Inches(stat_width - 0.3), Inches(0.8))
        tf = label_box.text_frame
        para = tf.paragraphs[0]
        para.text = stat.get("label", "")
        para.font.size = Pt(16)
        para.font.color.rgb = RgbColor(*hex_to_rgb(palette["text_dark"]))
        para.alignment = PP_ALIGN.CENTER

def _add_image_left_content(slide, content: SlideContent, palette: dict):
    """Image Left Layout"""
    if content.image_path:
        img_path = ASSETS_DIR / "images" / content.image_path
        if img_path.exists():
            slide.shapes.add_picture(str(img_path), Inches(0.5), Inches(1.6), width=Inches(5.5))
    
    # Text on right
    text_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.6), Inches(6), Inches(5))
    tf = text_box.text_frame
    tf.word_wrap = True
    if content.content:
        para = tf.paragraphs[0]
        para.text = content.content
        para.font.size = Pt(16)
        para.font.color.rgb = RgbColor(*hex_to_rgb(palette["text_dark"]))

def _add_image_right_content(slide, content: SlideContent, palette: dict):
    """Image Right Layout"""
    # Text on left
    text_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(6), Inches(5))
    tf = text_box.text_frame
    tf.word_wrap = True
    if content.content:
        para = tf.paragraphs[0]
        para.text = content.content
        para.font.size = Pt(16)
        para.font.color.rgb = RgbColor(*hex_to_rgb(palette["text_dark"]))
    
    if content.image_path:
        img_path = ASSETS_DIR / "images" / content.image_path
        if img_path.exists():
            slide.shapes.add_picture(str(img_path), Inches(7.3), Inches(1.6), width=Inches(5.5))

# ============================================================================
# WORD DOCUMENT CREATION
# ============================================================================

def create_professional_docx(request: WordDocumentRequest) -> Path:
    """Erstellt ein professionelles Word-Dokument"""
    
    # Template laden oder neues Dokument
    template_path = get_template_path(request.template, "docx")
    if template_path:
        doc = Document(template_path)
    else:
        doc = Document()
    
    # Logo hinzufügen
    logo_path = get_logo_path(request.logo)
    if logo_path:
        doc.add_picture(str(logo_path), width=DocxInches(2))
        doc.add_paragraph()
    
    # Titel
    title_para = doc.add_heading(request.title, level=0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # Subtitle
    if request.subtitle:
        subtitle_para = doc.add_paragraph(request.subtitle)
        subtitle_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in subtitle_para.runs:
            run.font.size = DocxPt(14)
            run.font.italic = True
    
    # Author and Date
    if request.author:
        meta_para = doc.add_paragraph(f"{request.author} | {datetime.now().strftime('%d.%m.%Y')}")
        meta_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    doc.add_paragraph()  # Spacer
    
    # Inhaltsverzeichnis Platzhalter
    if request.include_toc:
        doc.add_heading("Inhaltsverzeichnis", level=1)
        doc.add_paragraph("[Inhaltsverzeichnis wird nach dem Öffnen in Word aktualisiert]")
        doc.add_page_break()
    
    # Sections
    for section in request.sections:
        if section.heading:
            doc.add_heading(section.heading, level=section.heading_level)
        
        if section.content:
            doc.add_paragraph(section.content)
        
        if section.bullet_points:
            for point in section.bullet_points:
                para = doc.add_paragraph(point, style='List Bullet')
        
        if section.table:
            _add_table_to_doc(doc, section.table)
        
        if section.image_path:
            img_path = ASSETS_DIR / "images" / section.image_path
            if img_path.exists():
                doc.add_picture(str(img_path), width=DocxInches(5))
    
    # Header/Footer
    if request.header_text or request.footer_text:
        section = doc.sections[0]
        if request.header_text:
            header = section.header
            header_para = header.paragraphs[0]
            header_para.text = request.header_text
            header_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        
        if request.footer_text:
            footer = section.footer
            footer_para = footer.paragraphs[0]
            footer_para.text = request.footer_text
            footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # Save
    filename = generate_filename("document", "docx")
    filepath = FILES_DIR / filename
    doc.save(filepath)
    return filepath

def _add_table_to_doc(doc, table_data: List[List[str]]):
    """Fügt eine formatierte Tabelle zum Dokument hinzu"""
    if not table_data:
        return
    
    table = doc.add_table(rows=len(table_data), cols=len(table_data[0]))
    table.style = 'Table Grid'
    
    for i, row in enumerate(table_data):
        for j, cell in enumerate(row):
            table.rows[i].cells[j].text = str(cell)
            # Header row styling
            if i == 0:
                for paragraph in table.rows[i].cells[j].paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True

# ============================================================================
# EXCEL CREATION (COMPLEX)
# ============================================================================

def create_professional_xlsx(request: ExcelRequest) -> Path:
    """Erstellt eine professionelle Excel-Datei mit Charts und Formatierung"""
    
    # Template laden oder neues Workbook
    template_path = get_template_path(request.template, "xlsx")
    if template_path:
        wb = load_workbook(template_path)
    else:
        wb = Workbook()
        # Default sheet entfernen wenn wir Sheets definieren
        if request.sheets and 'Sheet' in wb.sheetnames:
            del wb['Sheet']
    
    # Styling
    header_font = Font(bold=True, color="FFFFFF", size=11)
    header_fill = PatternFill(start_color="003366", end_color="003366", fill_type="solid")
    
    if request.style == "financial":
        header_fill = PatternFill(start_color="1E3A5F", end_color="1E3A5F", fill_type="solid")
    elif request.style == "minimal":
        header_fill = PatternFill(start_color="666666", end_color="666666", fill_type="solid")
    
    border = Border(
        left=Side(style='thin'),
        right=Side(style='thin'),
        top=Side(style='thin'),
        bottom=Side(style='thin')
    )
    
    for sheet_config in request.sheets:
        if sheet_config.name in wb.sheetnames:
            ws = wb[sheet_config.name]
        else:
            ws = wb.create_sheet(title=sheet_config.name)
        
        current_row = 1
        
        # Headers
        if sheet_config.headers:
            for col, header in enumerate(sheet_config.headers, 1):
                cell = ws.cell(row=current_row, column=col, value=header)
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = Alignment(horizontal='center', vertical='center')
                cell.border = border
            current_row += 1
        
        # Data
        if sheet_config.data:
            for row_data in sheet_config.data:
                for col, value in enumerate(row_data, 1):
                    cell = ws.cell(row=current_row, column=col, value=value)
                    cell.border = border
                    cell.alignment = Alignment(horizontal='center' if isinstance(value, (int, float)) else 'left')
                current_row += 1
        
        # Formulas
        if sheet_config.formulas:
            for cell_ref, formula in sheet_config.formulas.items():
                ws[cell_ref] = formula
                ws[cell_ref].font = Font(bold=True)
        
        # Column widths
        if sheet_config.column_widths:
            for col, width in sheet_config.column_widths.items():
                ws.column_dimensions[col].width = width
        else:
            # Auto-width
            for col in range(1, ws.max_column + 1):
                ws.column_dimensions[get_column_letter(col)].width = 15
        
        # Charts
        if sheet_config.charts:
            for chart_config in sheet_config.charts:
                chart = _create_chart(ws, chart_config)
                if chart:
                    ws.add_chart(chart, chart_config.position)
    
    # Logo
    logo_path = get_logo_path(request.logo)
    if logo_path:
        img = XLImage(str(logo_path))
        img.width = 150
        img.height = 50
        first_sheet = wb.worksheets[0]
        first_sheet.add_image(img, 'A1')
    
    # Save
    filename = generate_filename(request.filename.replace('.xlsx', ''), "xlsx")
    filepath = FILES_DIR / filename
    wb.save(filepath)
    return filepath

def _create_chart(ws, config: ChartConfig):
    """Erstellt einen Chart basierend auf der Konfiguration"""
    
    # Parse data range
    parts = config.data_range.split(':')
    if len(parts) != 2:
        return None
    
    # Simplified chart creation
    if config.type == "bar":
        chart = BarChart()
    elif config.type == "line":
        chart = LineChart()
    elif config.type == "pie":
        chart = PieChart()
    else:
        chart = BarChart()
    
    chart.title = config.title
    chart.style = 10
    
    # Note: Full data reference setup would require parsing the range
    # This is a simplified version
    
    return chart

# ============================================================================
# PDF CREATION
# ============================================================================

def create_professional_pdf(request: PDFRequest) -> Path:
    """Erstellt ein professionelles PDF-Dokument"""
    
    filename = generate_filename("document", "pdf")
    filepath = FILES_DIR / filename
    
    page_size = A4 if request.page_size == "A4" else letter
    
    doc = SimpleDocTemplate(
        str(filepath),
        pagesize=page_size,
        rightMargin=2*cm,
        leftMargin=2*cm,
        topMargin=2*cm,
        bottomMargin=2*cm
    )
    
    styles = getSampleStyleSheet()
    palette = COLOR_PALETTES.get(request.color_scheme, COLOR_PALETTES["corporate_blue"])
    
    # Custom styles
    styles.add(ParagraphStyle(
        name='CustomTitle',
        parent=styles['Title'],
        fontSize=24,
        textColor=colors.HexColor(f"#{palette['primary']}"),
        spaceAfter=30
    ))
    
    styles.add(ParagraphStyle(
        name='CustomHeading1',
        parent=styles['Heading1'],
        fontSize=18,
        textColor=colors.HexColor(f"#{palette['primary']}"),
        spaceBefore=20,
        spaceAfter=10
    ))
    
    styles.add(ParagraphStyle(
        name='CustomHeading2',
        parent=styles['Heading2'],
        fontSize=14,
        textColor=colors.HexColor(f"#{palette['secondary']}"),
        spaceBefore=15,
        spaceAfter=8
    ))
    
    story = []
    
    # Logo
    logo_path = get_logo_path(request.logo)
    if logo_path:
        img = RLImage(str(logo_path), width=2*inch, height=0.75*inch)
        story.append(img)
        story.append(Spacer(1, 0.5*inch))
    
    # Title
    story.append(Paragraph(request.title, styles['CustomTitle']))
    
    # Author
    if request.author:
        story.append(Paragraph(f"<i>{request.author} | {datetime.now().strftime('%d.%m.%Y')}</i>", styles['Normal']))
        story.append(Spacer(1, 0.5*inch))
    
    # Sections
    for section in request.sections:
        if section.type == "heading":
            style = styles['CustomHeading1'] if section.level == 1 else styles['CustomHeading2']
            story.append(Paragraph(section.content or "", style))
        
        elif section.type == "paragraph":
            story.append(Paragraph(section.content or "", styles['Normal']))
            story.append(Spacer(1, 0.2*inch))
        
        elif section.type == "table" and section.data:
            table = Table(section.data)
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor(f"#{palette['primary']}")),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 12),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('TEXTCOLOR', (0, 1), (-1, -1), colors.black),
                ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
                ('FONTSIZE', (0, 1), (-1, -1), 10),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
            ]))
            story.append(table)
            story.append(Spacer(1, 0.3*inch))
        
        elif section.type == "image" and section.image_path:
            img_path = ASSETS_DIR / "images" / section.image_path
            if img_path.exists():
                img = RLImage(str(img_path), width=4*inch, height=3*inch)
                story.append(img)
                story.append(Spacer(1, 0.3*inch))
        
        elif section.type == "spacer":
            height = section.height or 1
            story.append(Spacer(1, height*cm))
        
        elif section.type == "page_break":
            story.append(PageBreak())
    
    doc.build(story)
    return filepath

# ============================================================================
# FASTAPI APP
# ============================================================================

@asynccontextmanager
async def lifespan(app: FastAPI):
    # Startup
    task = asyncio.create_task(cleanup_old_files())
    yield
    # Shutdown
    task.cancel()

app = FastAPI(
    title="File Creator Service - Extended",
    description="Erstellt professionelle PPTX, DOCX, XLSX und PDF Dateien mit Design-Optionen, Templates und Logo-Support",
    version="2.0.0",
    lifespan=lifespan
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ============================================================================
# API ENDPOINTS
# ============================================================================

@app.get("/")
async def root():
    return {
        "service": "File Creator Service - Extended",
        "version": "2.0.0",
        "endpoints": {
            "create_pptx": "/create/pptx",
            "create_docx": "/create/docx",
            "create_xlsx": "/create/xlsx",
            "create_pdf": "/create/pdf",
            "upload_logo": "/assets/logo",
            "upload_template": "/templates/upload",
            "list_assets": "/assets",
            "list_templates": "/templates",
            "download": "/files/{filename}",
            "list_files": "/files",
            "color_palettes": "/palettes"
        }
    }

# --- PowerPoint ---
@app.post("/create/pptx")
async def create_pptx(request: PresentationRequest):
    try:
        filepath = create_professional_pptx(request)
        return {
            "success": True,
            "filename": filepath.name,
            "download_url": f"{BASE_URL}/files/{filepath.name}",
            "message": f"PowerPoint '{request.title}' wurde erstellt."
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# --- Word Document ---
@app.post("/create/docx")
async def create_docx(request: WordDocumentRequest):
    try:
        filepath = create_professional_docx(request)
        return {
            "success": True,
            "filename": filepath.name,
            "download_url": f"{BASE_URL}/files/{filepath.name}",
            "message": f"Word-Dokument '{request.title}' wurde erstellt."
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# --- Excel ---
@app.post("/create/xlsx")
async def create_xlsx(request: ExcelRequest):
    try:
        filepath = create_professional_xlsx(request)
        return {
            "success": True,
            "filename": filepath.name,
            "download_url": f"{BASE_URL}/files/{filepath.name}",
            "message": f"Excel-Datei wurde erstellt."
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# --- PDF ---
@app.post("/create/pdf")
async def create_pdf(request: PDFRequest):
    try:
        filepath = create_professional_pdf(request)
        return {
            "success": True,
            "filename": filepath.name,
            "download_url": f"{BASE_URL}/files/{filepath.name}",
            "message": f"PDF '{request.title}' wurde erstellt."
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# --- Asset Management ---
@app.post("/assets/logo")
async def upload_logo(file: UploadFile = File(...), name: Optional[str] = Form(None)):
    try:
        logo_name = name or file.filename
        logo_path = ASSETS_DIR / "logos" / logo_name
        
        with open(logo_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        
        return LogoUploadResponse(
            success=True,
            logo_name=logo_name,
            message=f"Logo '{logo_name}' wurde hochgeladen."
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/assets/image")
async def upload_image(file: UploadFile = File(...), name: Optional[str] = Form(None)):
    try:
        image_name = name or file.filename
        image_path = ASSETS_DIR / "images" / image_name
        
        with open(image_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        
        return {
            "success": True,
            "image_name": image_name,
            "message": f"Bild '{image_name}' wurde hochgeladen."
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/assets")
async def list_assets():
    logos = list((ASSETS_DIR / "logos").glob("*"))
    images = list((ASSETS_DIR / "images").glob("*"))
    
    return {
        "logos": [f.name for f in logos if f.is_file()],
        "images": [f.name for f in images if f.is_file()]
    }

# --- Template Management ---
@app.post("/templates/upload")
async def upload_template(
    file: UploadFile = File(...),
    template_type: str = Form(...),  # docx, xlsx, pptx
    name: Optional[str] = Form(None)
):
    try:
        template_name = name or file.filename
        template_dir = TEMPLATES_DIR / template_type
        template_dir.mkdir(exist_ok=True)
        template_path = template_dir / template_name
        
        with open(template_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        
        return TemplateUploadResponse(
            success=True,
            template_name=template_name,
            template_type=template_type,
            message=f"Template '{template_name}' wurde hochgeladen."
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/templates")
async def list_templates():
    templates = {}
    for template_type in ["docx", "xlsx", "pptx"]:
        type_dir = TEMPLATES_DIR / template_type
        if type_dir.exists():
            templates[template_type] = [f.name for f in type_dir.glob("*") if f.is_file()]
        else:
            templates[template_type] = []
    return templates

# --- Color Palettes ---
@app.get("/palettes")
async def get_palettes():
    return {
        "available_palettes": list(COLOR_PALETTES.keys()),
        "palettes": COLOR_PALETTES
    }

# --- File Management ---
@app.get("/files/{filename}")
async def download_file(filename: str):
    filepath = FILES_DIR / filename
    if not filepath.exists():
        raise HTTPException(status_code=404, detail="Datei nicht gefunden")
    return FileResponse(filepath, filename=filename)

@app.get("/files")
async def list_files():
    files = []
    for f in FILES_DIR.iterdir():
        if f.is_file():
            files.append({
                "filename": f.name,
                "download_url": f"{BASE_URL}/files/{f.name}",
                "size_bytes": f.stat().st_size,
                "created": datetime.fromtimestamp(f.stat().st_mtime).isoformat()
            })
    return {"files": files}

@app.delete("/files/{filename}")
async def delete_file(filename: str):
    filepath = FILES_DIR / filename
    if not filepath.exists():
        raise HTTPException(status_code=404, detail="Datei nicht gefunden")
    filepath.unlink()
    return {"success": True, "message": f"Datei '{filename}' wurde gelöscht."}

# ============================================================================
# MAIN
# ============================================================================

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=PORT)
